#!/usr/bin/env python2
# -*- coding: utf-8 -*-
"""
Created on Thu Oct 12 16:14:01 2017

@author: yxcheng
"""


#%%
from pptx import Presentation
prs = Presentation()
#blank_slide_layout = prs.slide_layouts[6]
#slide = prs.slides.add_slide(blank_slide_layout)
#img_path = './example/Picture4.png'
#add_pic_with_title(slide, img_path, 5, 4.5, 2.5, 'Picture4.png')
#prs.save('test2.pptx')

import os
img_paths = []
for file in os.listdir("./example"):
    if file.endswith(".png"):
        print(os.path.join("./example", file))
        img_paths.append(os.path.join("./example", file))



img_name = ['0','1','2','3','4','5','6','7']
img_nums = 6
prs = add_new_slide_pics(prs, img_paths, img_nums, img_name)
prs.save('test1.pptx')
#%%
def add_pic_with_title(slide, img_path, left, top, width, img_name):
    # add a picture to slide from img_path and add the title with img_name
    from pptx.util import Inches
    left = Inches(left)
    top_img = Inches(top)
    width = Inches(width)
    
    # add picture
    pic = slide.shapes.add_picture(img_path, left, top_img, width=width)
    
    #add title
    top_title = Inches(top - 0.3)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top_title , width, height)
    tf = txBox.text_frame
    tf.text = img_name
    
    return slide
    
def add_new_slide_pics(prs, img_paths, img_nums, img_name):
#    from pptx.util import Inches
    #img_nums = len(img_paths)

    #%% 2
    if img_nums == 2:
        layout = [1, 2]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)
    
    #%% 3
    elif img_nums == 3:
        layout = [1, 3]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)
   
    #%% 4
    elif img_nums == 4:
        layout = [2, 2]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)
    
    #%% 6
    elif img_nums == 6:
        layout = [2, 3]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)
    
    #%% 8
    elif img_nums == 8:
        layout = [2, 4]
        prs = set_layout(prs, layout, img_paths, img_nums, img_name)        
               
    else:
        print('other numbers')
    return prs
    
def set_layout(prs, layout,img_paths, img_nums, img_name):
    #decide the layout of picture by accepting the layout_num
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    ind = 0
    int_width = 0.3
    width = (10 - int_width)/layout[1]
    pic_width = width - int_width
    int_top = 7.5 - layout[0]*width
#    for img_path in img_paths: 
    for i in range(0, img_nums):
        img_path = img_paths[i]          
        col = ind%layout[1]
        row = ind//layout[1]
#            height = ((10 - int_width)/layout[1] - int_width
        left = col*width + int_width
        top = row*width + int_top
        
        # add the picture on its position
        add_pic_with_title(slide, img_path, left, top, pic_width, img_name[ind])
        ind = ind + 1
    return prs
    
    
######################################################
#%% test part
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')

#%%
from pptx import Presentation

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p. level = 2

prs.save('test.pptx')
#%%
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = "This is text inside a textbox"

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold"
p.font.bold = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big"
p.font.size = Pt(40)

prs.save('test.pptx')
#%%
from pptx import Presentation
from pptx.util import Inches



prs = Presentation() 
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

img_path = './example/Picture1.png'

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]


img_path = './example/Picture2.png'
left = Inches(5)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]
#%%
img_path = './example/Picture2.png'
left = Inches(5)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top*0.7 , width, height)
tf = txBox.text_frame

tf.text = img_path[10::]
    
    
    